import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import numpy as np
import os
import re
from datetime import datetime, timedelta
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load environment variables from .env file
import load_env

# Thomson Reuters brand colors
TR_ORANGE = '#D64000'
TR_RACING_GREEN = '#123015'
TR_WHITE = '#FFFFFF'

# Secondary colors - Sky pair
TR_LIGHT_SKY = '#E3F1FD'
TR_DARK_SKY = '#0874E3'

# Additional secondary colors
TR_LIGHT_GOLD = '#FCF2DA'
TR_DARK_GOLD = '#E9B045'
TR_LIGHT_AMBER = '#F8EADD'
TR_DARK_AMBER = '#D4792A'
TR_LIGHT_TEAL = '#E3F3EE'
TR_DARK_TEAL = '#4DB299'
TR_LIGHT_LIME = '#E1F4CD'
TR_DARK_LIME = '#8FCB64'

# Create a custom color palette for charts
TR_COLORS = [TR_ORANGE, TR_RACING_GREEN, TR_DARK_SKY, TR_DARK_AMBER, 
             TR_DARK_TEAL, TR_DARK_GOLD, TR_DARK_LIME]

# Extended color palette for pie charts to avoid color duplication
TR_PIE_COLORS = [
    TR_ORANGE,          # Primary orange
    TR_RACING_GREEN,    # Primary racing green
    TR_DARK_SKY,        # Dark sky blue
    TR_DARK_AMBER,      # Dark amber
    TR_DARK_TEAL,       # Dark teal
    TR_DARK_GOLD,       # Dark gold
    TR_DARK_LIME,       # Dark lime
    '#6c757d',          # Gray
    '#9C27B0',          # Purple
    '#FF5722',          # Deep orange
    '#607D8B',          # Blue gray
    '#795548',          # Brown
    '#009688',          # Teal
    '#673AB7',          # Deep purple
    '#3F51B5',          # Indigo
    '#2196F3',          # Blue
    '#00BCD4',          # Cyan
    '#4CAF50',          # Green
    '#8BC34A',          # Light green
    '#CDDC39',          # Lime
    '#FFEB3B',          # Yellow
    '#FFC107',          # Amber
    '#FF9800',          # Orange
    '#9E9E9E',          # Gray
    '#E91E63'           # Pink
]

# Set style for matplotlib
plt.style.use('ggplot')
plt.rcParams['axes.prop_cycle'] = plt.cycler(color=TR_COLORS)
plt.rcParams['figure.facecolor'] = TR_WHITE
plt.rcParams['axes.facecolor'] = TR_WHITE
plt.rcParams['axes.edgecolor'] = TR_RACING_GREEN
plt.rcParams['axes.labelcolor'] = TR_RACING_GREEN
plt.rcParams['xtick.color'] = TR_RACING_GREEN
plt.rcParams['ytick.color'] = TR_RACING_GREEN
plt.rcParams['grid.color'] = TR_LIGHT_SKY
plt.rcParams['text.color'] = TR_RACING_GREEN
plt.rcParams['font.family'] = ['Clario', 'Arial', 'sans-serif']

# Set font sizes according to Reuters typography guidelines
plt.rcParams['font.size'] = 10  # Base font size
plt.rcParams['axes.titlesize'] = 20  # Headline 2 (3X base)
plt.rcParams['axes.labelsize'] = 12  # 1.2X base
plt.rcParams['xtick.labelsize'] = 10  # Base size
plt.rcParams['ytick.labelsize'] = 10  # Base size
plt.rcParams['legend.fontsize'] = 10  # Base size
plt.rcParams['figure.titlesize'] = 30  # Headline 1 (6X base) - for figure suptitles

# Create output directories if they don't exist
os.makedirs('output', exist_ok=True)
os.makedirs('output/images', exist_ok=True)

# Load the data
print("Loading data...")
# Try to load the actual data file, if not available, use the sample data
try:
    file_path = 'Sample_Data/TV5 test_2025_07_23_18_19_36.xlsx'
    df = pd.read_excel(file_path)
    print(f"Loaded actual Teletrax data from {file_path}")
except FileNotFoundError:
    try:
        file_path = 'Sample_Data/Sample_Teletrax_Data.xlsx'
        df = pd.read_excel(file_path)
        print(f"Loaded sample data from {file_path}")
    except FileNotFoundError:
        print("No data file found. Please run create_sample_data.py to generate sample data or place your Teletrax export in the Sample_Data directory.")
        exit(1)

# Convert datetime columns to datetime type
datetime_cols = ['UTC detection start', 'Local detection start', 
                'Hit: Activation date start (UTC)', 'Asset: Activation date start (UTC)']
for col in datetime_cols:
    df[col] = pd.to_datetime(df[col])

# Convert 'Actual detection length' to seconds for easier analysis
def timedelta_to_seconds(td):
    if isinstance(td, timedelta):
        return td.total_seconds()
    elif isinstance(td, (int, float)):
        return float(td)
    else:
        return 0.0

df['Detection Length (seconds)'] = df['Actual detection length'].apply(timedelta_to_seconds)

# Add derived columns for analysis
df['Detection Year'] = df['UTC detection start'].dt.year
df['Detection Month'] = df['UTC detection start'].dt.month
df['Detection Day'] = df['UTC detection start'].dt.day
df['Detection Hour'] = df['UTC detection start'].dt.hour
df['Detection Weekday'] = df['UTC detection start'].dt.day_name()
df['Detection Date'] = df['UTC detection start'].dt.date

# Extract topic from slug line (text before the first slash)
df['Topic'] = df['Slug line'].apply(lambda x: x.split('-')[0] if '-' in x else x.split('/')[0])
df['Subtopic'] = df['Slug line'].apply(lambda x: x.split('-')[1].split('/')[0] if '-' in x else '')

print(f"Data loaded. {len(df)} records spanning from {df['UTC detection start'].min().date()} to {df['UTC detection start'].max().date()}")

# List of countries for extraction
COUNTRIES = [
    'Afghanistan', 'Albania', 'Algeria', 'Andorra', 'Angola', 'Argentina', 'Armenia', 'Australia', 
    'Austria', 'Azerbaijan', 'Bahamas', 'Bahrain', 'Bangladesh', 'Barbados', 'Belarus', 'Belgium', 
    'Belize', 'Benin', 'Bhutan', 'Bolivia', 'Bosnia', 'Botswana', 'Brazil', 'Brunei', 'Bulgaria', 
    'Burkina Faso', 'Burundi', 'Cambodia', 'Cameroon', 'Canada', 'Cape Verde', 'Central African Republic', 
    'Chad', 'Chile', 'China', 'Colombia', 'Comoros', 'Congo', 'Costa Rica', 'Croatia', 'Cuba', 'Cyprus', 
    'Czech Republic', 'Denmark', 'Djibouti', 'Dominica', 'Dominican Republic', 'East Timor', 'Ecuador', 
    'Egypt', 'El Salvador', 'Equatorial Guinea', 'Eritrea', 'Estonia', 'Eswatini', 'Ethiopia', 'Fiji', 
    'Finland', 'France', 'Gabon', 'Gambia', 'Georgia', 'Germany', 'Ghana', 'Greece', 'Grenada', 
    'Guatemala', 'Guinea', 'Guinea-Bissau', 'Guyana', 'Haiti', 'Honduras', 'Hungary', 'Iceland', 'India', 
    'Indonesia', 'Iran', 'Iraq', 'Ireland', 'Israel', 'Italy', 'Ivory Coast', 'Jamaica', 'Japan', 
    'Jordan', 'Kazakhstan', 'Kenya', 'Kiribati', 'Kosovo', 'Kuwait', 'Kyrgyzstan', 'Laos', 'Latvia', 
    'Lebanon', 'Lesotho', 'Liberia', 'Libya', 'Liechtenstein', 'Lithuania', 'Luxembourg', 'Madagascar', 
    'Malawi', 'Malaysia', 'Maldives', 'Mali', 'Malta', 'Marshall Islands', 'Mauritania', 'Mauritius', 
    'Mexico', 'Micronesia', 'Moldova', 'Monaco', 'Mongolia', 'Montenegro', 'Morocco', 'Mozambique', 
    'Myanmar', 'Namibia', 'Nauru', 'Nepal', 'Netherlands', 'New Zealand', 'Nicaragua', 'Niger', 'Nigeria', 
    'North Korea', 'North Macedonia', 'Norway', 'Oman', 'Pakistan', 'Palau', 'Palestine', 'Panama', 
    'Papua New Guinea', 'Paraguay', 'Peru', 'Philippines', 'Poland', 'Portugal', 'Qatar', 'Romania', 
    'Russia', 'Rwanda', 'Saint Kitts and Nevis', 'Saint Lucia', 'Saint Vincent and the Grenadines', 
    'Samoa', 'San Marino', 'Sao Tome and Principe', 'Saudi Arabia', 'Senegal', 'Serbia', 'Seychelles', 
    'Sierra Leone', 'Singapore', 'Slovakia', 'Slovenia', 'Solomon Islands', 'Somalia', 'South Africa', 
    'South Korea', 'South Sudan', 'Spain', 'Sri Lanka', 'Sudan', 'Suriname', 'Sweden', 'Switzerland', 
    'Syria', 'Taiwan', 'Tajikistan', 'Tanzania', 'Thailand', 'Togo', 'Tonga', 'Trinidad and Tobago', 
    'Tunisia', 'Turkey', 'Turkmenistan', 'Tuvalu', 'Uganda', 'Ukraine', 'United Arab Emirates', 
    'United Kingdom', 'United States', 'Uruguay', 'Uzbekistan', 'Vanuatu', 'Vatican City', 'Venezuela', 
    'Vietnam', 'Yemen', 'Zambia', 'Zimbabwe'
]

# Add common country abbreviations and alternative names
COUNTRY_ALIASES = {
    'USA': 'United States', 'US': 'United States', 'America': 'United States', 'U.S.': 'United States',
    'UK': 'United Kingdom', 'Britain': 'United Kingdom', 'Great Britain': 'United Kingdom',
    'UAE': 'United Arab Emirates',
    'N Korea': 'North Korea', 'S Korea': 'South Korea',
    'DRC': 'Congo', 'Democratic Republic of Congo': 'Congo',
    'Ivory Coast': 'Côte d\'Ivoire',
    'Russia': 'Russian Federation',
    'Czech': 'Czech Republic',
    'UAE': 'United Arab Emirates',
    'Saudi': 'Saudi Arabia',
    'Palestine': 'Palestinian Territories',
    'Bosnia': 'Bosnia and Herzegovina',
    'Myanmar': 'Burma',
    'Eswatini': 'Swaziland'
}

def extract_countries_from_text(text):
    """Extract country names from text
    
    This function extracts country names from text by looking for matches in the
    COUNTRIES list and COUNTRY_ALIASES dictionary. It returns a list of unique
    country names found in the text.
    
    Args:
        text (str): Text to extract countries from
        
    Returns:
        list: List of unique country names found in the text
    """
    if not isinstance(text, str):
        return []
    
    text = text.strip()
    countries_found = []
    
    # Check for countries in the text
    for country in COUNTRIES:
        if country in text:
            countries_found.append(country)
    
    # Check for country aliases
    for alias, country in COUNTRY_ALIASES.items():
        if alias in text and country not in countries_found:
            countries_found.append(country)
    
    # If no countries found, check for first word as potential country
    if not countries_found and text:
        first_word = text.split()[0] if text.split() else ""
        if first_word in COUNTRIES or first_word in COUNTRY_ALIASES:
            if first_word in COUNTRY_ALIASES:
                countries_found.append(COUNTRY_ALIASES[first_word])
            else:
                countries_found.append(first_word)
    
    return list(set(countries_found))

# ===== ANALYSIS FUNCTIONS =====

def generate_time_series_analysis():
    """Generate time series analysis of detections"""
    print("Generating time series analysis...")
    
    # Aggregate by date
    daily_counts = df.groupby('Detection Date').size()
    
    # Resample to fill in missing dates with zeros
    date_range = pd.date_range(start=daily_counts.index.min(), end=daily_counts.index.max())
    daily_counts = daily_counts.reindex(date_range, fill_value=0)
    
    # Calculate 30-day moving average
    moving_avg = daily_counts.rolling(window=30).mean()
    
    # Plot time series
    plt.figure(figsize=(15, 7))
    plt.plot(daily_counts.index, daily_counts.values, color=TR_DARK_SKY, alpha=0.5, label='Daily Detections')
    plt.plot(moving_avg.index, moving_avg.values, color=TR_ORANGE, linewidth=2, label='30-Day Moving Average')
    plt.title('Reuters Video Detections Over Time', fontsize=16)
    plt.xlabel('Date', fontsize=12)
    plt.ylabel('Number of Detections', fontsize=12)
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Format x-axis to show years
    plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m'))
    plt.gca().xaxis.set_major_locator(mdates.YearLocator())
    plt.xticks(rotation=45)
    
    plt.tight_layout()
    plt.savefig('output/images/time_series_analysis.png', dpi=300)
    plt.close()
    
    # Create interactive plotly version
    fig = px.line(x=daily_counts.index, y=daily_counts.values, 
                 labels={'x': 'Date', 'y': 'Number of Detections'},
                 title='Reuters Video Detections Over Time',
                 color_discrete_sequence=[TR_DARK_SKY])
    fig.add_scatter(x=moving_avg.index, y=moving_avg.values, mode='lines', 
                   name='30-Day Moving Average', line=dict(color=TR_ORANGE, width=2))
    fig.update_layout(
        template='plotly_white', 
        legend_title_text='',
        paper_bgcolor=TR_WHITE,
        plot_bgcolor=TR_WHITE,
        font_color=TR_RACING_GREEN,
        title_font_color=TR_RACING_GREEN
    )
    fig.write_html('output/images/time_series_analysis_interactive.html')
    
    # Channel comparison over time
    channel_daily = df.groupby(['Detection Date', 'Channel: Name']).size().unstack(fill_value=0)
    
    # Resample to fill in missing dates
    channel_daily = channel_daily.reindex(date_range, fill_value=0)
    
    # Plot channel comparison
    plt.figure(figsize=(15, 7))
    for channel in channel_daily.columns:
        plt.plot(channel_daily.index, channel_daily[channel], 
                label=channel, alpha=0.7, linewidth=2)
    
    plt.title('Channel Comparison: Detections Over Time', fontsize=16)
    plt.xlabel('Date', fontsize=12)
    plt.ylabel('Number of Detections', fontsize=12)
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Format x-axis
    plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m'))
    plt.gca().xaxis.set_major_locator(mdates.YearLocator())
    plt.xticks(rotation=45)
    
    plt.tight_layout()
    plt.savefig('output/images/channel_comparison_time_series.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.line(channel_daily, labels={'value': 'Number of Detections', 'Detection Date': 'Date'},
                 title='Channel Comparison: Detections Over Time')
    fig.update_layout(template='plotly_white', legend_title_text='Channel')
    fig.write_html('output/images/channel_comparison_interactive.html')
    
    return daily_counts, channel_daily

def analyze_top_stories():
    """Analyze the most frequently detected stories
    
    This function creates visualizations of the most frequently detected stories,
    organized both by Story ID and by Slug Line. Stories are grouped by Story ID
    with the most common headline displayed for each ID.
    """
    print("Analyzing top stories...")
    
    # Group by Story ID and find the most common headline for each ID
    story_id_counts = df.groupby('Story ID').size().sort_values(ascending=False).head(15)
    
    # Function to clean headline text
    def clean_headline(text):
        if not isinstance(text, str):
            return "No headline"
        
        # Remove HTML tags
        import re
        clean_text = re.sub(r'<.*?>', '', text)
        
        # Remove special characters and normalize whitespace
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()
        
        # If the text is empty or just whitespace after cleaning, return "No headline"
        if not clean_text or clean_text.isspace():
            return "No headline"
        
        # Truncate long headlines to a reasonable length (60 characters)
        if len(clean_text) > 60:
            clean_text = clean_text[:57] + "..."
            
        return clean_text
    
    # For each Story ID, find the most common headline
    story_id_headlines = {}
    for story_id in story_id_counts.index:
        headlines = df[df['Story ID'] == story_id]['Headline'].value_counts()
        if not headlines.empty and pd.notna(headlines.index[0]):
            most_common_headline = clean_headline(headlines.index[0])
        else:
            most_common_headline = "No headline"
        story_id_headlines[story_id] = most_common_headline
    
    # Create labels with only the most common headline, but ensure uniqueness
    # If multiple Story IDs have the same headline, we'll display them separately
    labels = []
    headline_counter = {}
    
    for story_id in story_id_counts.index:
        headline = story_id_headlines[story_id]
        
        # If this headline has been seen before, add a counter to make it unique
        if headline in headline_counter:
            headline_counter[headline] += 1
            labels.append(f"{headline} ({headline_counter[headline]})")
        else:
            headline_counter[headline] = 1
            labels.append(headline)
    
    # Top stories by slug line (keep this for comparison)
    top_stories_by_slug = df['Slug line'].value_counts().head(15)
    
    # Plot top stories by ID with most common headline
    plt.figure(figsize=(14, 10))
    bars = plt.barh(labels[::-1], story_id_counts.values[::-1], color=TR_ORANGE)
    plt.title('Top 15 Most Detected Stories (by Story ID)', fontsize=16)
    plt.xlabel('Number of Detections', fontsize=12)
    plt.ylabel('Most Common Headline', fontsize=12)
    plt.grid(True, alpha=0.3, axis='x')
    
    # Set a fixed x-axis limit to prevent stretching
    # Calculate a reasonable maximum based on the data
    max_value = story_id_counts.max()
    plt.xlim(0, max_value * 1.2)  # Add 20% padding
    
    # Add explanatory note
    plt.figtext(0.5, 0.01, 
               "Note: Stories are grouped by Story ID. Each story is displayed with its most commonly used headline.", 
               ha='center', fontsize=10, style='italic')
    
    # Add count labels to bars
    for bar in bars:
        width = bar.get_width()
        plt.text(width + 5, bar.get_y() + bar.get_height()/2, 
                f'{width:,.0f}', ha='left', va='center', fontsize=10)
    
    plt.tight_layout()
    plt.savefig('output/images/top_stories_by_id.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(y=labels[::-1], x=story_id_counts.values[::-1], 
                orientation='h', text=story_id_counts.values[::-1],
                labels={'y': 'Most Common Headline', 'x': 'Number of Detections'},
                title='Top 15 Most Detected Stories (by Story ID)')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/top_stories_by_id_interactive.html')
    
    # Extract master slugs (portion before the forward slash)
    def extract_master_slug(slug_line):
        if pd.isna(slug_line):
            return "Unknown"
        
        # Replace "ADVISORY " prefix with "LIVE: " if present (ADVISORY indicates a live broadcast)
        if isinstance(slug_line, str) and slug_line.startswith("ADVISORY "):
            slug_line = "LIVE: " + slug_line[9:].strip()
        
        # Find the position of the first forward slash
        slash_pos = slug_line.find('/')
        
        # If a slash is found, extract everything before it and add the slash
        if slash_pos != -1:
            return slug_line[:slash_pos+1]  # Include the slash
        else:
            # If no slash is found, return the entire slug line
            return slug_line
    
    # Apply the function to extract master slugs
    df['Master Slug'] = df['Slug line'].apply(extract_master_slug)
    
    # Count occurrences of each master slug
    top_master_slugs = df['Master Slug'].value_counts().head(15)
    
    # Create the top themes visualization
    plt.figure(figsize=(12, 8))
    bars = plt.barh(top_master_slugs.index[::-1], top_master_slugs.values[::-1], color=TR_DARK_TEAL)
    plt.title('Top 15 Thematic Areas (by Master Slug)', fontsize=16)
    plt.xlabel('Number of Detections', fontsize=12)
    plt.ylabel('Master Slug', fontsize=12)
    plt.grid(True, alpha=0.3, axis='x')
    
    # Add count labels to bars
    for bar in bars:
        width = bar.get_width()
        plt.text(width + 5, bar.get_y() + bar.get_height()/2, 
                f'{width:,.0f}', ha='left', va='center', fontsize=10)
    
    plt.tight_layout()
    plt.savefig('output/images/top_themes.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(y=top_master_slugs.index[::-1], x=top_master_slugs.values[::-1], 
                orientation='h', text=top_master_slugs.values[::-1],
                labels={'y': 'Master Slug', 'x': 'Number of Detections'},
                title='Top 15 Thematic Areas (by Master Slug)')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/top_themes_interactive.html')
    
    # Top topics analysis
    top_topics = df['Topic'].value_counts().head(10)
    
    plt.figure(figsize=(12, 8))
    bars = plt.barh(top_topics.index[::-1], top_topics.values[::-1], color=TR_DARK_TEAL)
    plt.title('Top 10 Most Detected Topics', fontsize=16)
    plt.xlabel('Number of Detections', fontsize=12)
    plt.ylabel('Topic', fontsize=12)
    plt.grid(True, alpha=0.3, axis='x')
    
    # Add count labels
    for bar in bars:
        width = bar.get_width()
        plt.text(width + 5, bar.get_y() + bar.get_height()/2, 
                f'{width:,.0f}', ha='left', va='center', fontsize=10)
    
    plt.tight_layout()
    plt.savefig('output/images/top_topics.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(y=top_topics.index[::-1], x=top_topics.values[::-1], 
                orientation='h', text=top_topics.values[::-1],
                labels={'y': 'Topic', 'x': 'Number of Detections'},
                title='Top 10 Most Detected Topics')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/top_topics_interactive.html')
    
    # Channel preference for top stories
    top_5_stories = df['Slug line'].value_counts().head(5).index
    channel_story_counts = pd.crosstab(df['Slug line'], df['Channel: Name'])
    channel_story_counts = channel_story_counts.loc[top_5_stories]
    
    # Plot channel preference
    plt.figure(figsize=(12, 8))
    channel_story_counts.plot(kind='barh', stacked=True, figsize=(12, 8), 
                             color=TR_COLORS)
    plt.title('Channel Preference for Top 5 Stories', fontsize=16)
    plt.xlabel('Number of Detections', fontsize=12)
    plt.ylabel('Story Slug Line', fontsize=12)
    plt.grid(True, alpha=0.3, axis='x')
    plt.legend(title='Channel')
    
    plt.tight_layout()
    plt.savefig('output/images/channel_story_preference.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(channel_story_counts, orientation='h', barmode='stack',
                labels={'value': 'Number of Detections', 'Slug line': 'Story Slug Line'},
                title='Channel Preference for Top 5 Stories')
    fig.update_layout(template='plotly_white', legend_title_text='Channel')
    fig.write_html('output/images/channel_story_preference_interactive.html')
    
    return story_id_counts, top_topics, channel_story_counts

def analyze_detection_patterns():
    """Analyze patterns in detection timing"""
    print("Analyzing detection patterns...")
    
    # Hour of day analysis
    hourly_counts = df['Detection Hour'].value_counts().sort_index()
    
    plt.figure(figsize=(12, 6))
    plt.bar(hourly_counts.index, hourly_counts.values, color=TR_DARK_SKY, alpha=0.7)
    plt.title('Detections by Hour of Day (UTC)', fontsize=16)
    plt.xlabel('Hour of Day', fontsize=12)
    plt.ylabel('Number of Detections', fontsize=12)
    plt.xticks(range(0, 24))
    plt.grid(True, alpha=0.3, axis='y')
    
    plt.tight_layout()
    plt.savefig('output/images/hourly_detection_pattern.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(x=hourly_counts.index, y=hourly_counts.values,
                labels={'x': 'Hour of Day (UTC)', 'y': 'Number of Detections'},
                title='Detections by Hour of Day')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/hourly_detection_pattern_interactive.html')
    
    # Day of week analysis
    # Ensure proper ordering of days
    day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
    weekday_counts = df['Detection Weekday'].value_counts().reindex(day_order)
    
    plt.figure(figsize=(12, 6))
    plt.bar(weekday_counts.index, weekday_counts.values, color=TR_DARK_TEAL, alpha=0.7)
    plt.title('Detections by Day of Week', fontsize=16)
    plt.xlabel('Day of Week', fontsize=12)
    plt.ylabel('Number of Detections', fontsize=12)
    plt.xticks(rotation=45)
    plt.grid(True, alpha=0.3, axis='y')
    
    plt.tight_layout()
    plt.savefig('output/images/weekday_detection_pattern.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(x=weekday_counts.index, y=weekday_counts.values,
                labels={'x': 'Day of Week', 'y': 'Number of Detections'},
                title='Detections by Day of Week',
                category_orders={"x": day_order})
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/weekday_detection_pattern_interactive.html')
    
    # Heatmap of hour by day of week
    hour_day_counts = df.groupby(['Detection Weekday', 'Detection Hour']).size().unstack(fill_value=0)
    # Reindex to ensure proper order
    hour_day_counts = hour_day_counts.reindex(day_order)
    
    plt.figure(figsize=(15, 8))
    plt.imshow(hour_day_counts, cmap='viridis', aspect='auto')
    plt.colorbar(label='Number of Detections')
    plt.title('Heatmap of Detections by Day and Hour', fontsize=16)
    plt.xlabel('Hour of Day (UTC)', fontsize=12)
    plt.ylabel('Day of Week', fontsize=12)
    plt.xticks(range(24), range(24))
    plt.yticks(range(len(day_order)), day_order)
    
    plt.tight_layout()
    plt.savefig('output/images/day_hour_heatmap.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.imshow(hour_day_counts, 
                   labels=dict(x="Hour of Day (UTC)", y="Day of Week", color="Number of Detections"),
                   title="Heatmap of Detections by Day and Hour")
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/day_hour_heatmap_interactive.html')
    
    return hourly_counts, weekday_counts, hour_day_counts

def analyze_detection_lengths():
    """Analyze the distribution of detection lengths"""
    print("Analyzing detection lengths...")
    
    # Distribution of detection lengths
    plt.figure(figsize=(12, 6))
    plt.hist(df['Detection Length (seconds)'], bins=30, color=TR_ORANGE, alpha=0.7)
    plt.title('Distribution of Detection Lengths', fontsize=16)
    plt.xlabel('Detection Length (seconds)', fontsize=12)
    plt.ylabel('Frequency', fontsize=12)
    plt.grid(True, alpha=0.3)
    
    plt.tight_layout()
    plt.savefig('output/images/detection_length_distribution.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.histogram(df, x='Detection Length (seconds)', nbins=30,
                      labels={'Detection Length (seconds)': 'Detection Length (seconds)'},
                      title='Distribution of Detection Lengths')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/detection_length_distribution_interactive.html')
    
    # Detection length by channel
    plt.figure(figsize=(10, 6))
    for channel in df['Channel: Name'].unique():
        channel_data = df[df['Channel: Name'] == channel]['Detection Length (seconds)']
        plt.hist(channel_data, bins=20, alpha=0.5, label=channel)
    
    plt.title('Detection Length Distribution by Channel', fontsize=16)
    plt.xlabel('Detection Length (seconds)', fontsize=12)
    plt.ylabel('Frequency', fontsize=12)
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    plt.tight_layout()
    plt.savefig('output/images/detection_length_by_channel.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.histogram(df, x='Detection Length (seconds)', color='Channel: Name', nbins=20,
                      barmode='overlay', opacity=0.7,
                      labels={'Detection Length (seconds)': 'Detection Length (seconds)'},
                      title='Detection Length Distribution by Channel')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/detection_length_by_channel_interactive.html')
    
    # Average detection length by story (top 10 stories)
    top_10_stories = df['Slug line'].value_counts().head(10).index
    story_lengths = df[df['Slug line'].isin(top_10_stories)].groupby('Slug line')['Detection Length (seconds)'].mean().sort_values(ascending=False)
    
    plt.figure(figsize=(12, 8))
    bars = plt.barh(story_lengths.index[::-1], story_lengths.values[::-1], color=TR_DARK_AMBER)
    plt.title('Average Detection Length for Top 10 Stories', fontsize=16)
    plt.xlabel('Average Detection Length (seconds)', fontsize=12)
    plt.ylabel('Story Slug Line', fontsize=12)
    plt.grid(True, alpha=0.3, axis='x')
    
    # Add value labels
    for bar in bars:
        width = bar.get_width()
        plt.text(width + 0.1, bar.get_y() + bar.get_height()/2, 
                f'{width:.1f}s', ha='left', va='center', fontsize=10)
    
    plt.tight_layout()
    plt.savefig('output/images/avg_detection_length_by_story.png', dpi=300)
    plt.close()
    
    # Interactive version
    fig = px.bar(y=story_lengths.index[::-1], x=story_lengths.values[::-1], 
                orientation='h', text=[f'{x:.1f}s' for x in story_lengths.values[::-1]],
                labels={'y': 'Story Slug Line', 'x': 'Average Detection Length (seconds)'},
                title='Average Detection Length for Top 10 Stories')
    fig.update_layout(template='plotly_white')
    fig.write_html('output/images/avg_detection_length_by_story_interactive.html')
    
    return story_lengths

def create_powerpoint_presentation(output_dir='output'):
    """Create a PowerPoint presentation with the analysis results
    
    This function generates a comprehensive PowerPoint presentation containing
    all the visualizations and analyses. It includes slides for both Story ID + Headline
    and Slug Line organization of stories.
    
    Args:
        output_dir (str): Directory to save the PowerPoint presentation to
    """
    print(f"Creating PowerPoint presentation in {output_dir}...")
    
    prs = Presentation()
    
    # Add Thomson Reuters logo to the presentation
    logo_path = 'static/images/logos/tr_pri_logo_rgb_color.png'
    
    # Set default font for all text in the presentation to Arial
    # (Arial is used for PowerPoint as it's shareable content per Reuters guidelines)
    for slide_layout in prs.slide_layouts:
        for shape in slide_layout.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
    
    # Convert hex colors to RGB for PowerPoint
    tr_orange_rgb = RGBColor.from_string(TR_ORANGE[1:])
    tr_racing_green_rgb = RGBColor.from_string(TR_RACING_GREEN[1:])
    tr_dark_sky_rgb = RGBColor.from_string(TR_DARK_SKY[1:])
    tr_white_rgb = RGBColor.from_string(TR_WHITE[1:])
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Add logo to title slide
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(2)
    slide.shapes.add_picture(logo_path, left, top, width=width)
    
    title.text = "RVN Client Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = tr_orange_rgb
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(40)  # Headline 1 size (approx 6X base)
    title.text_frame.paragraphs[0].font.bold = True  # Medium weight equivalent
    
    # Get unique channel names from the data
    channel_names = ", ".join(df['Channel: Name'].unique())
    subtitle.text = f"Analysis of {channel_names} Video Usage Data\n{datetime.now().strftime('%B %d, %Y')}"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(20)  # Subhead size (2X base)
    
    # Add background color to the slide
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = tr_white_rgb
    
    # Overview slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Analysis Overview"
    title.text_frame.paragraphs[0].font.color.rgb = tr_orange_rgb
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(40)  # Headline 1 size
    title.text_frame.paragraphs[0].font.bold = True  # Medium weight equivalent
    
    tf = body.text_frame
    tf.text = "This presentation includes:"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(20)  # Subhead size
    tf.paragraphs[0].font.bold = True  # Medium weight equivalent
    
    # Add background color to the slide
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = tr_white_rgb
    
    p = tf.add_paragraph()
    p.text = "Time series analysis of video detections"
    p.level = 1
    p.font.name = 'Arial'
    p.font.size = Pt(16)  # Body text size
    
    p = tf.add_paragraph()
    p.text = "Top stories and topics analysis"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Detection pattern analysis (by hour, day)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Detection length analysis"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Channel comparison"
    p.level = 1
    
    # Function to add an image slide with adaptive layout
    def add_image_slide(title_text, img_path, description=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = tr_orange_rgb
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.size = Pt(32)  # Headline 2 size
        title.text_frame.paragraphs[0].font.bold = True  # Medium weight equivalent
        
        # Add background color to the slide
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = tr_white_rgb
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Calculate available space for content (excluding title)
        title_height = Inches(1.2)  # Approximate height of title
        content_height = slide_height - title_height
        
        # Use the full path with output_dir
        full_img_path = os.path.join(output_dir, img_path.replace('output/', ''))
        
        # Get image dimensions to maintain aspect ratio
        from PIL import Image
        try:
            with Image.open(full_img_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
        except:
            # Default aspect ratio if image can't be opened
            aspect_ratio = 16/9
        
        # Calculate image dimensions based on content type
        # Different chart types need different handling
        if 'pie' in img_path:
            # Pie charts often need more height than width
            img_width = min(Inches(6.5), slide_width * 0.7)
            img_height = img_width / aspect_ratio
        elif 'bar' in img_path or 'stories' in img_path:
            # Bar charts often need more width
            img_width = min(Inches(7.5), slide_width * 0.8)
            img_height = img_width / aspect_ratio
        else:
            # Default sizing for other charts
            img_width = min(Inches(7), slide_width * 0.75)
            img_height = img_width / aspect_ratio
        
        # Ensure image isn't too tall
        max_img_height = content_height * 0.7
        if img_height > max_img_height:
            img_height = max_img_height
            img_width = img_height * aspect_ratio
        
        # Center the image horizontally
        left = (slide_width - img_width) / 2
        # Position image below title with some padding
        top = title_height + Inches(0.2)
        
        # Add image with calculated dimensions
        slide.shapes.add_picture(full_img_path, left, top, width=img_width, height=img_height)
        
        # Add description if provided
        if description:
            # Position text below image with padding
            text_top = top + img_height + Inches(0.3)
            # Ensure text doesn't go off slide
            max_text_height = slide_height - text_top - Inches(0.5)
            
            # Create text box
            text_width = min(Inches(8), slide_width * 0.85)
            text_left = (slide_width - text_width) / 2
            
            txBox = slide.shapes.add_textbox(text_left, text_top, text_width, max_text_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = description
            
            # Apply consistent formatting
            for paragraph in tf.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(14)  # Body text size
                paragraph.font.color.rgb = tr_racing_green_rgb
    
    # Add time series analysis slides
    add_image_slide("Video Detections Over Time", 
                   "output/images/time_series_analysis.png",
                   "The chart shows the daily number of detections with a 30-day moving average.")
    
    # Get channel names for the description
    channel_list = ", ".join(df['Channel: Name'].unique())
    add_image_slide("Channel Comparison Over Time", 
                   "output/images/channel_comparison_time_series.png",
                   f"Comparison of detection counts between {channel_list} channels.")
    
    # Add top stories slides with improved explanations
    add_image_slide("Top 15 Most Detected Stories (by Story ID)", 
                   "output/images/top_stories_by_id.png",
                   "Stories grouped by Story ID, with each ID displayed with its most common headline. " +
                   "This view shows which specific content pieces were most frequently detected.")
    
    add_image_slide("Top 15 Most Detected Stories (by Slug Line)", 
                   "output/images/top_stories_by_slug.png",
                   "Stories grouped by Slug Line. Multiple Story IDs may share the same Slug Line, " +
                   "so this view shows which broader topics or events received the most coverage.")
    
    # Add geographic distribution slide
    # Extract countries from slug lines and headlines
    content_countries = []
    
    # Process each story in the data
    for _, row in df.iterrows():
        # Extract countries from slug line (handle NaN values)
        slug_countries = extract_countries_from_text(row['Slug line']) if pd.notna(row['Slug line']) else []
        
        # Extract countries from headline (handle NaN values)
        headline_countries = extract_countries_from_text(row['Headline']) if pd.notna(row['Headline']) else []
        
        # Combine countries from both sources
        story_countries = list(set(slug_countries + headline_countries))
        
        # If no countries found, try using the topic as a potential country
        if not story_countries and 'Topic' in row and pd.notna(row['Topic']):
            topic_countries = extract_countries_from_text(row['Topic'])
            if topic_countries:
                story_countries = topic_countries
        
        # Add to the overall list
        content_countries.extend(story_countries)
    
    # Count occurrences of each country
    if content_countries:
        country_counts = pd.Series(content_countries).value_counts()
        
        # If there are too many countries, group smaller ones as "Rest of the world"
        if len(country_counts) > 8:
            top_countries = country_counts.head(7)
            rest_of_world = pd.Series({'Rest of the world': country_counts[7:].sum()})
            country_counts = pd.concat([top_countries, rest_of_world])
        
        # Calculate percentages
        country_percentages = (country_counts / country_counts.sum() * 100).round(0).astype(int)
    else:
        # If no countries found, create a default "Unknown" category
        country_counts = pd.Series({'Unknown': 1})
        country_percentages = pd.Series({'Unknown': 100})
    
    # Create the pie chart
    plt.figure(figsize=(10, 8))
    plt.pie(country_percentages, 
            labels=[f"{c}\n{p}%" for c, p in zip(country_percentages.index, country_percentages.values)], 
            autopct='', 
            startangle=90, 
            wedgeprops={'edgecolor': 'white'},
            colors=TR_PIE_COLORS)
    plt.title("Reuters Video Content Origin by Country", fontsize=16)
    plt.axis('equal')
    
    # Save the chart
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'images', 'country_distribution_pie.png'), dpi=300)
    plt.close()
    
    add_image_slide("Geographic Distribution of Stories", 
                   "output/images/country_distribution_pie.png",
                   "Distribution of Reuters video content by country of origin. Countries are extracted from story slug lines and headlines.")
    
    add_image_slide("Top 10 Most Detected Topics", 
                   "output/images/top_topics.png",
                   "The most frequently detected topics extracted from slug lines.")
    
    add_image_slide("Channel Preference for Top Stories", 
                   "output/images/channel_story_preference.png",
                   "How different channels use the top 5 most detected stories.")
    
    # Add detection pattern slides
    add_image_slide("Detections by Hour of Day", 
                   "output/images/hourly_detection_pattern.png",
                   "Distribution of detections across hours of the day (UTC).")
    
    add_image_slide("Detections by Day of Week", 
                   "output/images/weekday_detection_pattern.png",
                   "Distribution of detections across days of the week.")
    
    add_image_slide("Heatmap of Detections by Day and Hour", 
                   "output/images/day_hour_heatmap.png",
                   "Heatmap showing detection patterns by day of week and hour of day.")
    
    # Add detection length slides
    add_image_slide("Distribution of Detection Lengths", 
                   "output/images/detection_length_distribution.png",
                   "Histogram showing the distribution of detection lengths in seconds.")
    
    add_image_slide("Detection Length by Channel", 
                   "output/images/detection_length_by_channel.png",
                   "Comparison of detection length distributions between channels.")
    
    add_image_slide("Average Detection Length for Top Stories", 
                   "output/images/avg_detection_length_by_story.png",
                   "Average detection length for the top 10 most detected stories.")
    
    # Conclusion slide with consistent formatting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Insights"
    title.text_frame.paragraphs[0].font.color.rgb = tr_orange_rgb
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(40)  # Headline 1 size
    title.text_frame.paragraphs[0].font.bold = True  # Medium weight equivalent
    
    tf = body.text_frame
    tf.text = "Summary of findings:"
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.size = Pt(20)  # Subhead size
    tf.paragraphs[0].font.bold = True  # Medium weight equivalent
    
    # Add background color to the slide
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = tr_white_rgb
    
    # Function to add a consistently formatted bullet point
    def add_bullet_point(text_frame, text, level=1):
        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = 'Arial'
        p.font.size = Pt(16)  # Body text size
        p.font.color.rgb = tr_racing_green_rgb
        # Add consistent bullet formatting
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        return p
    
    # These will be populated with actual insights after running the analysis
    add_bullet_point(tf, "Most detected stories are related to security and politics in African countries")
    add_bullet_point(tf, "Average detection length is around 5-6 seconds")
    add_bullet_point(tf, "Detection patterns show higher usage during weekdays")
    
    # Get channel names for the insight
    channel_list = list(df['Channel: Name'].unique())
    if len(channel_list) >= 2:
        add_bullet_point(tf, f"{channel_list[0]} and {channel_list[1]} show different preferences for certain stories")
    else:
        add_bullet_point(tf, f"{channel_list[0]} shows distinct preferences for different story types")
    
    # Save the presentation
    output_path = os.path.join(output_dir, 'RVN_Client_Analysis.pptx')
    prs.save(output_path)
    print(f"PowerPoint presentation saved to '{output_path}'")
    return output_path

def create_single_slide_presentation(channel_name, output_dir='output', client_friendly=False):
    """Create a single-slide PowerPoint presentation for a specific channel
    
    This function generates a simple single-slide PowerPoint presentation for a specific
    channel, following the format of the sample slides. The slide has a 2x2 grid layout:
    - Top-left: Time series chart showing video usage over time with a trend line
    - Top-right: Pie chart showing video usage by country
    - Bottom-left: Text summary with key statistics and trends
    - Bottom-right: Major stories/topics with percentages
    
    The layout is carefully designed to ensure all elements fit properly within their
    designated areas, with consistent formatting and proper spacing.
    
    Args:
        channel_name (str): Name of the channel to create the presentation for
        output_dir (str): Directory to save the PowerPoint presentation to
        client_friendly (bool): Whether to create a client-friendly version that focuses on positive trends
    """
    print(f"Creating single-slide PowerPoint presentation for {channel_name}...")
    
    # Filter data for the specified channel
    channel_df = df[df['Channel: Name'] == channel_name]
    
    if len(channel_df) == 0:
        print(f"No data found for channel {channel_name}")
        return
    
    # Create a new presentation
    prs = Presentation()
    
    # Add Thomson Reuters logo to the presentation
    logo_path = 'static/images/logos/tr_pri_logo_rgb_color.png'
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Set slide background to TR white
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(TR_WHITE[1:])
    
    # Add logo to slide
    left = Inches(0.1)
    top = Inches(0.1)
    width = Inches(1)
    slide.shapes.add_picture(logo_path, left, top, width=width)
    
    # Create the 2x2 grid layout
    # Define the grid dimensions
    left_margin = Inches(0.1)
    top_margin = Inches(0.1)
    slide_width = prs.slide_width - 2 * left_margin
    slide_height = prs.slide_height - 2 * top_margin
    cell_width = slide_width / 2
    cell_height = slide_height / 2
    
    # ===== TOP-LEFT: Time series chart =====
    # Generate time series data for this channel
    monthly_counts = channel_df.groupby(pd.Grouper(key='UTC detection start', freq='M')).size()
    
    # For client-friendly version, find periods of growth
    if client_friendly:
        # Try to use LiteLLM for time series analysis if available
        try:
            import json
            import requests
            from litellm_client import LiteLLMClient
            
            # Initialize the LiteLLM client
            api_key = os.environ.get('LITELLM_API_KEY')
            api_url = os.environ.get('LITELLM_API_URL', 'https://litellm.int.thomsonreuters.com')
            
            if not api_key:
                raise ValueError("LITELLM_API_KEY environment variable is not set. Please set it in your .env file.")
            
            litellm_client = LiteLLMClient(api_key=api_key, api_url=api_url)
            
            # Prepare data for LiteLLM
            # Convert numpy int64 values to regular Python integers for JSON serialization
            time_series_data = {
                "monthly_data": {str(date.date()): int(count) for date, count in zip(monthly_counts.index, monthly_counts.values)},
                "channel_name": channel_name,
                "total_months": len(monthly_counts),
                "most_recent_date": str(monthly_counts.index.max().date()),
                "earliest_date": str(monthly_counts.index.min().date())
            }
            
            # Create a custom prompt for time series analysis
            system_prompt = """
You are a data analysis expert specializing in time series analysis. Your task is to analyze Reuters video usage data and identify the optimal time period to display in a client presentation.

For client presentations, we want to show the longest possible timeframe that demonstrates an overall positive growth trend, while always including the most recent data. Some natural fluctuations are acceptable as long as the overall trend line is positive.

Analyze the monthly data and return a JSON object with:
1. start_date: The recommended start date for the time series chart (YYYY-MM-DD format)
2. end_date: The end date, which should always be the most recent date in the dataset (YYYY-MM-DD format)
3. explanation: A brief explanation of why this timeframe was selected
"""
            
            user_prompt = f"""
Please analyze the following monthly video usage data for {channel_name} and identify the optimal timeframe to display in a client presentation.

Monthly data: {json.dumps(time_series_data["monthly_data"], indent=2)}

Most recent date: {time_series_data["most_recent_date"]}
Earliest date: {time_series_data["earliest_date"]}
Total months of data: {time_series_data["total_months"]}

Remember:
1. Always include the most recent data (up to {time_series_data["most_recent_date"]})
2. Find the longest timeframe that still shows an overall positive trend
3. Some fluctuations are acceptable as long as the trend line is positive
4. Return your analysis as a JSON object with start_date, end_date, and explanation fields
"""
            
            # Make the API request
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "anthropic/claude-sonnet-4-20250514",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.2,
                "max_tokens": 1000,
                "response_format": {"type": "json_object"}
            }
            
            response = requests.post(
                f"{api_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=30
            )
            
            response.raise_for_status()
            result = response.json()
            
            # Extract the recommended timeframe
            ai_response = json.loads(result['choices'][0]['message']['content'])
            start_date = pd.to_datetime(ai_response['start_date'])
            end_date = pd.to_datetime(ai_response['end_date'])
            
            # Filter the monthly counts to the recommended timeframe
            monthly_counts = monthly_counts[(monthly_counts.index >= start_date) & (monthly_counts.index <= end_date)]
            
            print(f"LiteLLM recommended timeframe: {start_date.date()} to {end_date.date()}")
            print(f"Explanation: {ai_response['explanation']}")
            
        except Exception as e:
            print(f"Error using LiteLLM for time series analysis: {str(e)}")
            # Fall back to traditional approach if LiteLLM fails
            print(f"Original data range for {channel_name}: {monthly_counts.index.min()} to {monthly_counts.index.max()}")
            
            # Try different approaches to find growth periods
            growth_found = False
            
            # Approach 1: Always include the most recent data and find the longest period with a positive trend
            # Start with the most recent month and work backwards
            end_date = monthly_counts.index.max()
            best_start_date = end_date
            best_slope = 0
            
            # Try different starting points, always ending at the most recent date
            for i in range(len(monthly_counts)):
                start_idx = len(monthly_counts) - i - 1
                if start_idx < 0:
                    break
                    
                # Get the data from this starting point to the end
                period_data = monthly_counts.iloc[start_idx:]
                
                # Calculate the trend line
                if len(period_data) >= 2:  # Need at least 2 points for a trend
                    z = np.polyfit(range(len(period_data)), period_data.values, 1)
                    
                    # If this period has a positive trend and is longer than our current best
                    if z[0] > 0 and (z[0] > best_slope or len(period_data) > len(monthly_counts.iloc[monthly_counts.index >= best_start_date])):
                        best_slope = z[0]
                        best_start_date = period_data.index[0]
                        growth_found = True
            
            if growth_found:
                # Filter to the best period found
                monthly_counts = monthly_counts[monthly_counts.index >= best_start_date]
                print(f"Found growth period using approach 1: {monthly_counts.index.min()} to {monthly_counts.index.max()} with slope {best_slope}")
            else:
                # If no growth period found, just use the most recent 6 months
                window_size = min(6, len(monthly_counts))
                monthly_counts = monthly_counts.iloc[-window_size:]
                print(f"No growth period found, using most recent {window_size} months: {monthly_counts.index.min()} to {monthly_counts.index.max()}")
    
    # Create the time series chart
    plt.figure(figsize=(5, 3))
    plt.plot(monthly_counts.index, monthly_counts.values, color=TR_DARK_SKY, linewidth=2)
    
    # Explicitly set the x-axis limits to only show the selected period
    plt.xlim(monthly_counts.index.min(), monthly_counts.index.max())
    
    # Add trend line with appropriate color based on trend direction
    z = np.polyfit(range(len(monthly_counts)), monthly_counts.values, 1)
    p = np.poly1d(z)
    
    # Determine if trend is positive or negative
    trend_is_positive = z[0] > 0
    
    # For client-friendly version, always show positive trend
    if client_friendly:
        # We've already selected a growth period, so just make sure the trend line is positive
        trend_is_positive = True
        
        # Recalculate the trend line for the selected period
        z = np.polyfit(range(len(monthly_counts)), monthly_counts.values, 1)
        p = np.poly1d(z)
    
    # Choose color based on trend direction
    trend_color = TR_DARK_TEAL if trend_is_positive else TR_ORANGE
    
    # For client-friendly version, always use a positive color
    if client_friendly:
        trend_color = TR_DARK_TEAL
    
    plt.plot(monthly_counts.index, p(range(len(monthly_counts))), color=trend_color, linewidth=2)
    
    # Format the chart
    title_text = f"{channel_name}: Reuters video assets per month"
    if client_friendly:
        title_text = f"{channel_name}: Recent growth in Reuters video usage"
    plt.title(title_text, fontsize=12)
    plt.grid(True, alpha=0.3)
    plt.xticks(rotation=90, fontsize=8)
    plt.yticks(fontsize=8)
    
    # Save the chart
    time_series_path = os.path.join(output_dir, f'{channel_name}_time_series.png')
    plt.tight_layout()
    plt.savefig(time_series_path, dpi=300)
    plt.close()
    
    # Add the chart to the slide
    left = left_margin
    top = top_margin
    width = cell_width
    height = cell_height
    slide.shapes.add_picture(time_series_path, left, top, width=width, height=height)
    
    # ===== TOP-RIGHT: Pie chart of content origin countries =====
    # Extract countries from slug lines and headlines
    content_countries = []
    
    # Process each story in the channel's data
    for _, row in channel_df.iterrows():
        # Extract countries from slug line (handle NaN values)
        slug_countries = extract_countries_from_text(row['Slug line']) if pd.notna(row['Slug line']) else []
        
        # Extract countries from headline (handle NaN values)
        headline_countries = extract_countries_from_text(row['Headline']) if pd.notna(row['Headline']) else []
        
        # Combine countries from both sources
        story_countries = list(set(slug_countries + headline_countries))
        
        # If no countries found, try using the topic as a potential country
        if not story_countries and 'Topic' in row and pd.notna(row['Topic']):
            topic_countries = extract_countries_from_text(row['Topic'])
            if topic_countries:
                story_countries = topic_countries
        
        # Add to the overall list
        content_countries.extend(story_countries)
    
    # Count occurrences of each country
    if content_countries:
        country_counts = pd.Series(content_countries).value_counts()
        
        # If there are too many countries, group smaller ones as "Rest of the world"
        if len(country_counts) > 8:
            top_countries = country_counts.head(7)
            rest_of_world = pd.Series({'Rest of the world': country_counts[7:].sum()})
            country_counts = pd.concat([top_countries, rest_of_world])
        
        # Calculate percentages
        country_percentages = (country_counts / country_counts.sum() * 100).round(0).astype(int)
    else:
        # If no countries found, create a default "Unknown" category
        country_counts = pd.Series({'Unknown': 1})
        country_percentages = pd.Series({'Unknown': 100})
    
    # Create the pie chart with optimized size and layout
    plt.figure(figsize=(5, 3))
    plt.pie(country_percentages, 
            labels=[f"{c}\n{p}%" for c, p in zip(country_percentages.index, country_percentages.values)], 
            autopct='', 
            startangle=90, 
            wedgeprops={'edgecolor': 'white'},
            colors=TR_PIE_COLORS,
            textprops={'fontsize': 8})  # Smaller font for labels to prevent overlap
    plt.title("Reuters video content origin by country", fontsize=12)
    plt.axis('equal')
    # Add tight layout to ensure pie fits within bounds
    plt.tight_layout(pad=0.5)
    
    # Save the chart
    pie_chart_path = os.path.join(output_dir, f'{channel_name}_country_pie.png')
    plt.tight_layout()
    plt.savefig(pie_chart_path, dpi=300)
    plt.close()
    
    # Add the chart to the slide
    left = left_margin + cell_width
    top = top_margin
    width = cell_width
    height = cell_height
    slide.shapes.add_picture(pie_chart_path, left, top, width=width, height=height)
    
    # ===== BOTTOM-LEFT: Text summary =====
    left = left_margin
    top = top_margin + cell_height
    width = cell_width
    height = cell_height
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Calculate trend information
    first_year = monthly_counts.index.min().year
    last_year = monthly_counts.index.max().year
    years_span = last_year - first_year
    
    # Calculate if usage has increased/decreased and by how much
    first_year_avg = channel_df[channel_df['Detection Year'] == first_year].shape[0] / 12
    last_year_avg = channel_df[channel_df['Detection Year'] == last_year].shape[0] / 12
    
    if last_year_avg > first_year_avg:
        trend_direction = "upwards"
        if last_year_avg > 2 * first_year_avg:
            trend_magnitude = "double"
        else:
            trend_magnitude = "gentle"
    else:
        trend_direction = "downwards"
        trend_magnitude = "gentle"
    
    # Count unique countries
    num_countries = channel_df['Market: Name'].nunique()
    
    # Count total detections
    total_detections = len(channel_df)
    
    # Add the text
    p = tf.add_paragraph()
    
    if client_friendly:
        # For client-friendly version, try to use LiteLLM for text generation
        try:
            import json
            import requests
            from litellm_client import LiteLLMClient
            
            # Initialize the LiteLLM client
            api_key = os.environ.get('LITELLM_API_KEY')
            api_url = os.environ.get('LITELLM_API_URL', 'https://litellm.int.thomsonreuters.com')
            
            if not api_key:
                raise ValueError("LITELLM_API_KEY environment variable is not set. Please set it in your .env file.")
            
            # Prepare data for LiteLLM
            # Calculate some key statistics for the prompt
            recent_months = channel_df[channel_df['Detection Year'] == last_year]
            recent_month_counts = recent_months.groupby(pd.Grouper(key='UTC detection start', freq='M')).size()
            
            # Check for growth
            growth_found = False
            growth_percent = 0
            
            # Check last month vs previous month
            if len(recent_month_counts) >= 2 and recent_month_counts.iloc[-1] > recent_month_counts.iloc[-2]:
                growth_found = True
                growth_percent = ((recent_month_counts.iloc[-1] / recent_month_counts.iloc[-2]) - 1) * 100
            
            # Top countries
            top_countries = channel_df['Market: Name'].value_counts().head(3).index.tolist()
            
            # Top topics
            top_topics = channel_df['Topic'].value_counts().head(3).index.tolist()
            
            # Prepare the data for the prompt
            # Convert numpy int64 values to regular Python integers for JSON serialization
            channel_data = {
                "channel_name": channel_name,
                "total_detections": int(total_detections),
                "num_countries": int(num_countries),
                "top_countries": top_countries,
                "top_topics": top_topics,
                "first_year": int(first_year),
                "last_year": int(last_year),
                "years_span": int(years_span),
                "trend_direction": trend_direction,
                "trend_magnitude": trend_magnitude,
                "growth_found": growth_found,
                "growth_percent": round(growth_percent) if growth_found else 0,
                "monthly_data": {str(date.date()): int(count) for date, count in zip(monthly_counts.index, monthly_counts.values)},
                "time_period": f"{monthly_counts.index.min().date()} to {monthly_counts.index.max().date()}"
            }
            
            # Create a custom prompt for text generation
            system_prompt = """
You are a data analyst at Reuters specializing in client engagement. Your task is to generate positive, client-friendly text for a PowerPoint slide about a client's usage of Reuters video content.

The text should be:
1. Positive and upbeat, highlighting growth and partnership
2. Concise (2-3 sentences maximum for each paragraph)
3. Professional but conversational
4. Focused on the client's success and the value of Reuters content

Generate two short paragraphs:
1. First paragraph: Highlight growth trends or partnership value
2. Second paragraph: Emphasize global reach and content diversity

Return a JSON object with:
1. paragraph1: The first paragraph text
2. paragraph2: The second paragraph text
"""
            
            user_prompt = f"""
Please generate client-friendly text for a PowerPoint slide about {channel_name}'s usage of Reuters video content.

Here are the key statistics:
- Channel: {channel_name}
- Total detections: {total_detections}
- Number of countries covered: {num_countries}
- Top countries: {', '.join(top_countries)}
- Top topics: {', '.join(top_topics)}
- Time period: {channel_data['time_period']}
- Trend direction: {trend_direction}
- Growth found: {growth_found}
- Growth percentage: {channel_data['growth_percent']}% (if applicable)

Remember to be positive and highlight the partnership value, even if the growth metrics aren't strong.
"""
            
            # Make the API request
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "anthropic/claude-sonnet-4-20250514",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.7,  # Higher temperature for more creative text
                "max_tokens": 500,
                "response_format": {"type": "json_object"}
            }
            
            response = requests.post(
                f"{api_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=30
            )
            
            response.raise_for_status()
            result = response.json()
            
            # Extract the generated text
            ai_response = json.loads(result['choices'][0]['message']['content'])
            
            # Use the AI-generated text
            p.text = ai_response['paragraph1']
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = ai_response['paragraph2']
            
            print(f"Used LiteLLM for text generation for {channel_name}")
            
        except Exception as e:
            print(f"Error using LiteLLM for text generation: {str(e)}")
            # Fall back to traditional approach if LiteLLM fails
            
            # Find the most recent month with growth
            recent_months = channel_df[channel_df['Detection Year'] == last_year]
            recent_month_counts = recent_months.groupby(pd.Grouper(key='UTC detection start', freq='M')).size()
            
            # Look for any growth in recent months
            growth_found = False
            growth_percent = 0
            
            # Check last month vs previous month
            if len(recent_month_counts) >= 2 and recent_month_counts.iloc[-1] > recent_month_counts.iloc[-2]:
                growth_found = True
                growth_percent = ((recent_month_counts.iloc[-1] / recent_month_counts.iloc[-2]) - 1) * 100
            
            # If no month-to-month growth, look for any 3-month period with growth
            if not growth_found and len(recent_month_counts) >= 3:
                for i in range(len(recent_month_counts) - 3):
                    if recent_month_counts.iloc[i+3] > recent_month_counts.iloc[i]:
                        growth_found = True
                        growth_percent = ((recent_month_counts.iloc[i+3] / recent_month_counts.iloc[i]) - 1) * 100
                        break
            
            if growth_found:
                # Round growth percentage to nearest whole number for cleaner presentation
                growth_percent_rounded = round(growth_percent)
                p.text = f"{channel_name} has shown strong growth in Reuters video usage, with a {growth_percent_rounded}% increase in the most recent month."
            else:
                # No recent growth, focus on overall relationship
                p.text = f"{channel_name} continues to be a valued partner for Reuters video content, consistently featuring our coverage across {num_countries} countries."
            
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.bold = True
            
            p = tf.add_paragraph()
            
            # Find the countries with the most coverage
            top_countries = channel_df['Market: Name'].value_counts().head(3).index.tolist()
            top_countries_str = ", ".join(top_countries)
            p.text = f"Their coverage prominently features Reuters video from {top_countries_str} and other key regions, demonstrating the global reach of our partnership."
    else:
        # Original text for internal version
        p.text = f"{channel_name}'s use of Reuters video progressed its {trend_magnitude} {trend_direction} trend in {last_year}, which has seen baseline usage {trend_magnitude} in the last {years_span} years."
        
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"It remains very reliant on Reuters video for coverage from the US and used nearly {total_detections:,} video assets from {num_countries} countries."
    
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    
    # ===== BOTTOM-RIGHT: Major stories =====
    left = left_margin + cell_width
    top = top_margin + cell_height
    width = cell_width
    height = cell_height
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Get top topics
    topic_counts = channel_df['Topic'].value_counts()
    total_topics = topic_counts.sum()
    
    # Calculate percentages
    topic_percentages = (topic_counts / total_topics * 100).round(1)
    
    # Get top 3 topics
    top_topics = topic_percentages.head(3)
    
    # Add the title
    p = tf.add_paragraph()
    p.text = "Major stories"
    p.font.name = 'Arial'
    p.font.size = Pt(36)  # Headline 1 size
    p.font.bold = True  # Medium weight equivalent
    # Convert hex color to RGB
    tr_orange_rgb = RGBColor.from_string(TR_ORANGE[1:])  # Remove the # from the hex color
    p.font.color.rgb = tr_orange_rgb
    
    # Add each topic with percentage
    for topic, percentage in top_topics.items():
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {topic}: "
        run.font.name = 'Arial'
        run.font.size = Pt(24)  # Headline 2 size
        run.font.bold = True
        
        run = p.add_run()
        run.text = f"{percentage}%"
        run.font.name = 'Arial'
        run.font.size = Pt(24)  # Headline 2 size
        run.font.color.rgb = tr_orange_rgb
    
    # Add footer text
    p = tf.add_paragraph()
    p.text = f"Of all detections of Reuters video on {channel_name} in {last_year}."
    p.font.name = 'Arial'
    p.font.size = Pt(12)  # Body text size
    
    # Save the presentation
    if client_friendly:
        output_path = os.path.join(output_dir, f'{channel_name}_Client_Friendly_Single_Slide.pptx')
    else:
        output_path = os.path.join(output_dir, f'{channel_name}_Single_Slide.pptx')
    
    prs.save(output_path)
    print(f"Single-slide PowerPoint saved to '{output_path}'")
    
    return output_path

def main():
    """Run all analyses and create PowerPoint presentations"""
    print("Running all analyses...")
    
    # Run all analyses
    daily_counts, channel_daily = generate_time_series_analysis()
    top_stories, top_topics, channel_story_counts = analyze_top_stories()
    hourly_counts, weekday_counts, hour_day_counts = analyze_detection_patterns()
    story_lengths = analyze_detection_lengths()
    
    # Create PowerPoint presentation
    create_powerpoint_presentation()
    
    # Create single-slide presentations for each channel
    for channel in df['Channel: Name'].unique():
        create_single_slide_presentation(channel)
    
    print("All analyses completed.")

if __name__ == "__main__":
    main()
